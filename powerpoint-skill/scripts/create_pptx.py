#!/usr/bin/env python3
"""
create_pptx.py - Main presentation builder with Zendesk branding

This script creates PowerPoint presentations following Zendesk Garden design system guidelines.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from typing import Dict, List, Optional

# Zendesk Garden Color Palette
COLORS = {
    # Primary Colors
    'blue_600': RGBColor(38, 148, 214),      # #2694d6 - Primary brand
    'blue_700': RGBColor(31, 115, 183),      # #1f73b7 - Interactive
    'grey_700': RGBColor(92, 105, 112),      # #5c6970 - Body text
    'grey_300': RGBColor(216, 220, 222),     # #d8dcde - Backgrounds
    'grey_100': RGBColor(248, 249, 249),     # #f8f9f9 - Light bg
    'green_600': RGBColor(38, 161, 120),     # #26a178 - Success
    'red_600': RGBColor(235, 92, 105),       # #eb5c69 - Danger
    'yellow_600': RGBColor(214, 115, 5),     # #d67305 - Warning

    # Secondary Colors
    'purple_600': RGBColor(178, 118, 205),   # #b276cd
    'teal_600': RGBColor(42, 157, 143),      # #2a9d8f
    'orange_600': RGBColor(213, 116, 40),    # #d57428
    'mint_600': RGBColor(22, 162, 96),       # #16a260

    # Neutral
    'white': RGBColor(255, 255, 255),
    'black': RGBColor(10, 13, 14),           # #0a0d0e
}

# Typography sizes (in points)
FONT_SIZES = {
    'xxxl': 44,  # Display headlines
    'xxl': 36,   # Section headers
    'xl': 28,    # Slide titles
    'lg': 20,    # Subheadings
    'md': 14,    # Body text
    'sm': 11,    # Captions
}

def create_presentation(title: str = "Untitled Presentation") -> Presentation:
    """Create a new presentation with Zendesk branding."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 aspect ratio (1920x1080)
    prs.slide_height = Inches(7.5)
    return prs

def add_title_slide(prs: Presentation, title: str, subtitle: str = "") -> None:
    """
    Add a title slide with Zendesk branding.

    Args:
        prs: Presentation object
        title: Main title text
        subtitle: Subtitle text (optional)
    """
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['white']

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5),
        Inches(2.5),
        Inches(12.333),
        Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(FONT_SIZES['xxxl'])
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = COLORS['blue_700']
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(4.2),
            Inches(12.333),
            Inches(1)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_frame.paragraphs[0].font.size = Pt(FONT_SIZES['lg'])
        subtitle_frame.paragraphs[0].font.color.rgb = COLORS['grey_700']
        subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Accent line
    line = slide.shapes.add_shape(
        1,  # Line shape
        Inches(4),
        Inches(5.5),
        Inches(5.333),
        Inches(0.05)
    )
    line.line.color.rgb = COLORS['blue_600']
    line.line.width = Pt(3)

def add_content_slide(
    prs: Presentation,
    title: str,
    content: List[str],
    layout_type: str = "bullet"
) -> None:
    """
    Add a content slide with bullet points or text.

    Args:
        prs: Presentation object
        title: Slide title
        content: List of content items (bullet points or paragraphs)
        layout_type: "bullet" or "text"
    """
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['white']

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5),
        Inches(0.5),
        Inches(12.333),
        Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(FONT_SIZES['xxl'])
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = COLORS['blue_700']

    # Content
    content_box = slide.shapes.add_textbox(
        Inches(0.75),
        Inches(1.8),
        Inches(11.833),
        Inches(5)
    )
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    for i, item in enumerate(content):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(FONT_SIZES['md'])
        p.font.color.rgb = COLORS['grey_700']
        p.space_after = Pt(12)

        if layout_type == "bullet":
            p.level = 0

def add_two_column_slide(
    prs: Presentation,
    title: str,
    left_content: List[str],
    right_content: List[str]
) -> None:
    """
    Add a two-column layout slide.

    Args:
        prs: Presentation object
        title: Slide title
        left_content: Content for left column
        right_content: Content for right column
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['white']

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5),
        Inches(0.5),
        Inches(12.333),
        Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(FONT_SIZES['xxl'])
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = COLORS['blue_700']

    # Left column
    left_box = slide.shapes.add_textbox(
        Inches(0.75),
        Inches(1.8),
        Inches(5.5),
        Inches(5)
    )
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    for i, item in enumerate(left_content):
        if i > 0:
            left_frame.add_paragraph()
        p = left_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(FONT_SIZES['md'])
        p.font.color.rgb = COLORS['grey_700']
        p.level = 0

    # Right column
    right_box = slide.shapes.add_textbox(
        Inches(6.75),
        Inches(1.8),
        Inches(5.5),
        Inches(5)
    )
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    for i, item in enumerate(right_content):
        if i > 0:
            right_frame.add_paragraph()
        p = right_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(FONT_SIZES['md'])
        p.font.color.rgb = COLORS['grey_700']
        p.level = 0

def add_section_header_slide(
    prs: Presentation,
    title: str,
    subtitle: str = ""
) -> None:
    """
    Add a section header slide to divide presentation sections.

    Args:
        prs: Presentation object
        title: Section title
        subtitle: Optional subtitle
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background with accent color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['blue_600']

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5),
        Inches(2.5),
        Inches(12.333),
        Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(FONT_SIZES['xxxl'])
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = COLORS['white']
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(4.2),
            Inches(12.333),
            Inches(1)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_frame.paragraphs[0].font.size = Pt(FONT_SIZES['lg'])
        subtitle_frame.paragraphs[0].font.color.rgb = COLORS['white']
        subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_closing_slide(
    prs: Presentation,
    title: str = "Thank You",
    contact_info: Optional[List[str]] = None
) -> None:
    """
    Add a closing/thank you slide.

    Args:
        prs: Presentation object
        title: Closing message
        contact_info: List of contact details (optional)
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['grey_100']

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5),
        Inches(2.5),
        Inches(12.333),
        Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(FONT_SIZES['xxxl'])
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = COLORS['blue_700']
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Contact info
    if contact_info:
        contact_box = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(4.5),
            Inches(12.333),
            Inches(2)
        )
        contact_frame = contact_box.text_frame
        for i, info in enumerate(contact_info):
            if i > 0:
                contact_frame.add_paragraph()
            p = contact_frame.paragraphs[i]
            p.text = info
            p.font.size = Pt(FONT_SIZES['sm'])
            p.font.color.rgb = COLORS['grey_700']
            p.alignment = PP_ALIGN.CENTER

def save_presentation(prs: Presentation, filename: str) -> str:
    """
    Save the presentation to a file.

    Args:
        prs: Presentation object
        filename: Output filename (should end with .pptx)

    Returns:
        Full path to saved file
    """
    if not filename.endswith('.pptx'):
        filename += '.pptx'

    prs.save(filename)
    return filename

# Example usage
if __name__ == "__main__":
    # Create a sample presentation
    prs = create_presentation("Sample Zendesk Presentation")

    # Add title slide
    add_title_slide(
        prs,
        "Zendesk Product Overview",
        "Building better customer experiences"
    )

    # Add content slide
    add_content_slide(
        prs,
        "Key Features",
        [
            "AI-powered customer support",
            "Omnichannel messaging",
            "Advanced analytics and reporting",
            "Customizable workflows",
            "Enterprise-grade security"
        ]
    )

    # Add section header
    add_section_header_slide(
        prs,
        "Product Demo",
        "See it in action"
    )

    # Add closing slide
    add_closing_slide(
        prs,
        "Thank You",
        [
            "Questions?",
            "contact@zendesk.com",
            "www.zendesk.com"
        ]
    )

    # Save
    save_presentation(prs, "sample_presentation.pptx")
    print("Sample presentation created successfully!")
